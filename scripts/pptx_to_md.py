from pptx import Presentation
import os
from PIL import Image


home_dir = "/home/syost/git/Cvirg_Pmarinus_RNAseq"
# Load presentation
ppt_file = os.path.join(home_dir, "analysis_ideas.pptx")
prs = Presentation(ppt_file)

# Create markdown file
md_filename = os.path.join(home_dir, "analysis_ideas.md")
img_dir = os.path.join(home_dir, "images")
os.makedirs(img_dir, exist_ok=True)


def extract_table(shape):
    """Extract table data from a shape and return Markdown-formatted table."""
    table = shape.table
    md_table = []
    
    # Extract header row
    headers = [cell.text.strip() for cell in table.rows[0].cells]
    md_table.append("| " + " | ".join(headers) + " |")
    md_table.append("|" + " --- |" * len(headers))

    # Extract table rows
    for row in table.rows[1:]:
        row_text = [cell.text.strip() for cell in row.cells]
        md_table.append("| " + " | ".join(row_text) + " |")
    
    return "\n".join(md_table)

with open(md_filename, "w", encoding="utf-8") as md_file:
    for i, slide in enumerate(prs.slides):
        md_file.write(f"# Slide {i+1}\n\n")

        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                lines = shape.text.strip().split("\n")
                
                # Treat first line as a possible header
                if len(lines) > 1:
                    md_file.write(f"## {lines[0]}\n\n")
                    md_file.write("\n".join(lines[1:]) + "\n\n")
                else:
                    md_file.write(f"{shape.text.strip()}\n\n")

            # Extract and save tables
            if shape.has_table:
                md_file.write(extract_table(shape) + "\n\n")

            # Extract and save images
            if hasattr(shape, "image") and shape.image:
                img_bytes = shape.image.blob
                img_path = os.path.join(img_dir, f"slide_{i+1}_image.png")

                with open(img_path, "wb") as img_file:
                    img_file.write(img_bytes)

                md_file.write(f"![Slide {i+1} Image](./{img_path})\n\n")

print(f"Markdown file '{md_filename}' created with extracted images and text.")
